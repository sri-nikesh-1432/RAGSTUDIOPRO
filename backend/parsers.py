"""
RAG Studio Pro - File Parsers
Real file parsing for PDF, DOCX, XLSX, PPTX, CSV, HTML, JSON, XML, Images, ZIP.
"""

import os
import json
import zipfile
import hashlib
import mimetypes
from pathlib import Path
from typing import Dict, Any, Optional, Tuple
from datetime import datetime


def detect_file_type(file_path: str) -> str:
    """Detect file type from extension."""
    ext = Path(file_path).suffix.lower()
    type_map = {
        ".pdf": "pdf",
        ".txt": "text",
        ".md": "markdown",
        ".csv": "csv",
        ".tsv": "csv",
        ".docx": "docx",
        ".doc": "docx",
        ".pptx": "pptx",
        ".ppt": "pptx",
        ".xlsx": "xlsx",
        ".xls": "xlsx",
        ".html": "html",
        ".htm": "html",
        ".json": "json",
        ".jsonl": "json",
        ".xml": "xml",
        ".svg": "xml",
        ".png": "image",
        ".jpg": "image",
        ".jpeg": "image",
        ".gif": "image",
        ".bmp": "image",
        ".tiff": "image",
        ".webp": "image",
        ".zip": "zip",
    }
    return type_map.get(ext, "unknown")


def get_file_metadata(file_path: str) -> Dict[str, Any]:
    """Get basic file metadata."""
    stat = os.stat(file_path)
    return {
        "file_path": file_path,
        "file_name": os.path.basename(file_path),
        "file_size": stat.st_size,
        "file_size_human": _human_size(stat.st_size),
        "modified_at": datetime.fromtimestamp(stat.st_mtime).isoformat(),
        "created_at": datetime.fromtimestamp(stat.st_ctime).isoformat(),
        "extension": Path(file_path).suffix.lower(),
    }


def _human_size(size: int) -> str:
    """Convert bytes to human-readable string."""
    for unit in ["B", "KB", "MB", "GB"]:
        if size < 1024:
            return f"{size:.1f} {unit}"
        size /= 1024
    return f"{size:.1f} TB"


def _analyze_text(text: str) -> Dict[str, Any]:
    """Analyze extracted text for statistics."""
    words = text.split()
    paragraphs = [p for p in text.split("\n\n") if p.strip()]
    sentences = []
    for w in text.replace("!", ".").replace("?", ".").split("."):
        w = w.strip()
        if w:
            sentences.append(w)

    return {
        "characters": len(text),
        "words": len(words),
        "paragraphs": len(paragraphs),
        "sentences": len(sentences),
        "lines": len(text.split("\n")),
        "language": _detect_language(text),
    }


def _detect_language(text: str) -> str:
    """Simple language detection based on character patterns."""
    if not text:
        return "unknown"
    # Check for common scripts
    sample = text[:1000]
    latin_chars = sum(1 for c in sample if c.isascii() and c.isalpha())
    cjk_chars = sum(1 for c in sample if '\u4e00' <= c <= '\u9fff' or '\u3040' <= c <= '\u30ff')
    arabic_chars = sum(1 for c in sample if '\u0600' <= c <= '\u06ff')
    cyrillic_chars = sum(1 for c in sample if '\u0400' <= c <= '\u04ff')

    total_alpha = latin_chars + cjk_chars + arabic_chars + cyrillic_chars
    if total_alpha == 0:
        return "unknown"

    if cjk_chars / total_alpha > 0.3:
        return "Chinese/Japanese"
    if arabic_chars / total_alpha > 0.3:
        return "Arabic"
    if cyrillic_chars / total_alpha > 0.3:
        return "Russian/Cyrillic"
    if latin_chars / total_alpha > 0.5:
        return "English/Latin"
    return "mixed"


# ─── PDF Parser ────────────────────────────────────────────────────

def parse_pdf(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """Parse PDF using PyMuPDF."""
    try:
        import pymupdf
        doc = pymupdf.open(file_path)
        text_parts = []
        pages = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_text = page.get_text()
            text_parts.append(page_text)
            pages.append({
                "page_number": page_num + 1,
                "characters": len(page_text),
                "words": len(page_text.split()),
            })

        text = "\n\n".join(text_parts)
        metadata = {
            "pages": len(doc),
            "page_details": pages,
            "pdf_metadata": dict(doc.metadata) if doc.metadata else {},
            "is_encrypted": doc.is_encrypted,
        }
        doc.close()
        return text, metadata
    except ImportError:
        return "", {"error": "pymupdf not installed"}
    except Exception as e:
        return "", {"error": str(e)}


# ─── DOCX Parser ───────────────────────────────────────────────────

def parse_docx(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """Parse Word documents."""
    try:
        from docx import Document
        doc = Document(file_path)

        text_parts = []
        paragraphs = []
        tables = []

        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
                paragraphs.append({
                    "text": para.text[:100] + "..." if len(para.text) > 100 else para.text,
                    "style": para.style.name if para.style else "Normal",
                })

        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
                text_parts.append(" | ".join(row_data))
            tables.append(table_data)

        text = "\n\n".join(text_parts)
        metadata = {
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables),
            "table_details": [{"rows": len(t), "cols": len(t[0]) if t else 0} for t in tables],
            "sections": len(doc.sections),
        }
        return text, metadata
    except ImportError:
        return "", {"error": "python-docx not installed"}
    except Exception as e:
        return "", {"error": str(e)}


# ─── XLSX Parser ───────────────────────────────────────────────────

def parse_xlsx(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """Parse Excel files."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)

        text_parts = []
        sheet_details = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = []
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                row_count += 1
                row_str = " | ".join([str(c) if c is not None else "" for c in row])
                if row_str.strip(" |"):
                    rows_text.append(row_str)

            text_parts.append(f"Sheet: {sheet_name}\n" + "\n".join(rows_text))
            sheet_details.append({
                "name": sheet_name,
                "rows": row_count,
                "columns": ws.max_column or 0,
            })

        text = "\n\n".join(text_parts)
        metadata = {
            "sheets": len(wb.sheetnames),
            "sheet_names": wb.sheetnames,
            "sheet_details": sheet_details,
        }
        wb.close()
        return text, metadata
    except ImportError:
        return "", {"error": "openpyxl not installed"}
    except Exception as e:
        return "", {"error": str(e)}


# ─── PPTX Parser ───────────────────────────────────────────────────

def parse_pptx(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """Parse PowerPoint files."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)

        text_parts = []
        slide_details = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            text_parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))
            slide_details.append({
                "slide_number": slide_num,
                "text_shapes": len([s for s in slide.shapes if hasattr(s, "text")]),
                "images": len([s for s in slide.shapes if s.shape_type == 13]),
            })

        text = "\n\n".join(text_parts)
        metadata = {
            "slides": len(prs.slides),
            "slide_details": slide_details,
            "slide_width": str(prs.slide_width),
            "slide_height": str(prs.slide_height),
        }
        return text, metadata
    except ImportError:
        return "", {"error": "python-pptx not installed"}
    except Exception as e:
        return "", {"error": str(e)}


# ─── CSV Parser ────────────────────────────────────────────────────

def parse_csv(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """Parse CSV files."""
    try:
        import csv
        text_parts = []
        row_count = 0
        headers = []

        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            sample = f.read(8192)
            f.seek(0)

            # Detect delimiter
            sniffer = csv.Sniffer()
            try:
                dialect = sniffer.sniff(sample)
            except csv.Error:
                dialect = csv.excel

            reader = csv.reader(f, dialect)
            for i, row in enumerate(reader):
                row_count += 1
                if i == 0:
                    headers = row
                    text_parts.append("Headers: " + " | ".join(row))
                else:
                    text_parts.append(" | ".join(row))

        text = "\n".join(text_parts)
        metadata = {
            "rows": row_count,
            "columns": len(headers),
            "headers": headers,
            "delimiter": getattr(dialect, "delimiter", ","),
        }
        return text, metadata
    except Exception as e:
        return "", {"error": str(e)}


# ─── HTML Parser ───────────────────────────────────────────────────

def parse_html(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """Parse HTML files."""
    try:
        from bs4 import BeautifulSoup

        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()

        soup = BeautifulSoup(content, "lxml")

        # Remove script and style elements
        for element in soup(["script", "style", "nav", "footer", "header"]):
            element.decompose()

        title = soup.title.string if soup.title else ""
        text = soup.get_text(separator="\n", strip=True)

        # Count elements
        all_tags = [tag.name for tag in soup.find_all()]
        metadata = {
            "title": title,
            "total_elements": len(all_tags),
            "links": len(soup.find_all("a")),
            "images": len(soup.find_all("img")),
            "tables": len(soup.find_all("table")),
            "headings": {
                f"h{i}": len(soup.find_all(f"h{i}"))
                for i in range(1, 7)
                if soup.find_all(f"h{i}")
            },
        }
        return text, metadata
    except ImportError:
        return "", {"error": "beautifulsoup4 not installed"}
    except Exception as e:
        return "", {"error": str(e)}


# ─── JSON Parser ───────────────────────────────────────────────────

def parse_json(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """Parse JSON files."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        text = json.dumps(data, indent=2, ensure_ascii=False)

        metadata = {
            "type": type(data).__name__,
            "size_bytes": os.path.getsize(file_path),
        }
        if isinstance(data, list):
            metadata["items"] = len(data)
        elif isinstance(data, dict):
            metadata["keys"] = list(data.keys())[:50]
            metadata["depth"] = _json_depth(data)

        return text, metadata
    except Exception as e:
        return "", {"error": str(e)}


def _json_depth(obj, current=0):
    """Calculate max depth of JSON object."""
    if isinstance(obj, dict):
        if not obj:
            return current
        return max(_json_depth(v, current + 1) for v in obj.values())
    elif isinstance(obj, list):
        if not obj:
            return current
        return max(_json_depth(v, current + 1) for v in obj)
    return current


# ─── XML Parser ────────────────────────────────────────────────────

def parse_xml(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """Parse XML files."""
    try:
        import xml.etree.ElementTree as ET
        tree = ET.parse(file_path)
        root = tree.getroot()

        text_parts = []
        _xml_to_text(root, text_parts, depth=0)

        metadata = {
            "root_tag": root.tag,
            "total_elements": sum(1 for _ in root.iter()),
            "attributes": dict(root.attrib),
        }
        return "\n".join(text_parts), metadata
    except Exception as e:
        return "", {"error": str(e)}


def _xml_to_text(element, parts, depth=0):
    """Recursively extract text from XML."""
    indent = "  " * depth
    if element.text and element.text.strip():
        parts.append(f"{indent}{element.tag}: {element.text.strip()}")
    for child in element:
        _xml_to_text(child, parts, depth + 1)


# ─── Image Parser (OCR) ───────────────────────────────────────────

def parse_image(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """Parse images using OCR."""
    try:
        from PIL import Image
        img = Image.open(file_path)

        metadata = {
            "format": img.format,
            "mode": img.mode,
            "width": img.width,
            "height": img.height,
            "megapixels": round((img.width * img.height) / 1_000_000, 2),
        }

        # Try OCR
        try:
            import pytesseract
            text = pytesseract.image_to_string(img)
            metadata["ocr_performed"] = True
            return text, metadata
        except ImportError:
            metadata["ocr_performed"] = False
            metadata["note"] = "pytesseract not installed - OCR unavailable"
            return "", metadata
        except Exception as e:
            metadata["ocr_error"] = str(e)
            return "", metadata
    except ImportError:
        return "", {"error": "Pillow not installed"}
    except Exception as e:
        return "", {"error": str(e)}


# ─── ZIP Parser ────────────────────────────────────────────────────

def parse_zip(file_path: str) -> Tuple[str, Dict[str, Any]]:
    """Parse ZIP archives - extract and parse contained files."""
    try:
        extracted_files = []
        all_text = []

        with zipfile.ZipFile(file_path, "r") as zf:
            file_list = zf.namelist()
            extractable = [
                f for f in file_list
                if not f.endswith("/") and not f.startswith("__MACOSX")
            ]

            # Extract to temp directory and parse
            import tempfile
            with tempfile.TemporaryDirectory() as tmpdir:
                zf.extractall(tmpdir)
                for fname in extractable[:100]:  # Limit to 100 files
                    full_path = os.path.join(tmpdir, fname)
                    if os.path.isfile(full_path):
                        ext = Path(fname).suffix.lower()
                        if ext in [".txt", ".md", ".csv", ".json", ".html", ".xml"]:
                            try:
                                with open(full_path, "r", encoding="utf-8", errors="replace") as f:
                                    content = f.read()
                                all_text.append(f"--- {fname} ---\n{content}")
                                extracted_files.append(fname)
                            except Exception:
                                pass
                        elif ext == ".pdf":
                            text, _ = parse_pdf(full_path)
                            if text:
                                all_text.append(f"--- {fname} ---\n{text}")
                                extracted_files.append(fname)
                        elif ext == ".docx":
                            text, _ = parse_docx(full_path)
                            if text:
                                all_text.append(f"--- {fname} ---\n{text}")
                                extracted_files.append(fname)

        text = "\n\n".join(all_text)
        metadata = {
            "total_files": len(file_list),
            "extractable_files": len(extractable),
            "parsed_files": extracted_files,
            "archive_size": os.path.getsize(file_path),
        }
        return text, metadata
    except Exception as e:
        return "", {"error": str(e)}


# ─── Main Parse Function ──────────────────────────────────────────

def parse_file(file_path: str, extract_metadata: bool = True) -> Dict[str, Any]:
    """
    Parse any supported file type and return text + metadata.
    This is the main entry point for file parsing.
    """
    if not os.path.exists(file_path):
        return {
            "success": False,
            "error": f"File not found: {file_path}",
            "file_name": os.path.basename(file_path),
            "file_type": "unknown",
        }

    file_type = detect_file_type(file_path)
    file_meta = get_file_metadata(file_path) if extract_metadata else {}

    parser_map = {
        "pdf": parse_pdf,
        "docx": parse_docx,
        "xlsx": parse_xlsx,
        "pptx": parse_pptx,
        "csv": parse_csv,
        "text": lambda fp: (open(fp, "r", encoding="utf-8", errors="replace").read(), {}),
        "markdown": lambda fp: (open(fp, "r", encoding="utf-8", errors="replace").read(), {}),
        "html": parse_html,
        "json": parse_json,
        "xml": parse_xml,
        "image": parse_image,
        "zip": parse_zip,
    }

    parser = parser_map.get(file_type)
    if not parser:
        # Try reading as text
        try:
            text = open(file_path, "r", encoding="utf-8", errors="replace").read()
            return {
                "success": True,
                "text": text,
                "file_name": os.path.basename(file_path),
                "file_type": file_type,
                "metadata": file_meta,
                **_analyze_text(text),
            }
        except Exception as e:
            return {
                "success": False,
                "error": f"Unsupported file type: {file_type}. Error: {e}",
                "file_name": os.path.basename(file_path),
                "file_type": file_type,
            }

    text, type_metadata = parser(file_path)
    text_stats = _analyze_text(text) if text else {}

    return {
        "success": bool(text),
        "text": text,
        "file_name": os.path.basename(file_path),
        "file_type": file_type,
        "file_size": file_meta.get("file_size", 0),
        "metadata": {**file_meta, **type_metadata},
        **text_stats,
    }
